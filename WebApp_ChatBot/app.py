# app.py
import os
import io
import re
import json
import uuid
from typing import List, Dict, Any

import streamlit as st
import pandas as pd
import numpy as np

from dotenv import load_dotenv
load_dotenv()

# (optional) set your working dir if you need it
os.chdir('/Users/jameslim/Desktop/GITHUB/ai_km_poc/WebApp_ChatBot')

# ========== Config ==========
OPENAI_API_KEY = os.getenv("OPENAI_API_KEY", "")
if not OPENAI_API_KEY:
    st.stop()

VECTOR_DB_PATH = os.getenv("VECTOR_DB_PATH", ".chroma")
SHAREPOINT_ENABLED = os.getenv("SHAREPOINT_ENABLED", "false").lower() == "true"

# SharePoint (certificate-based app-only)
SP_TENANT      = os.getenv("SP_TENANT")
SP_SITE_URL    = os.getenv("SP_SITE_URL")
SP_LIST_NAME   = os.getenv("SP_LIST_NAME")
SP_CLIENT_ID   = os.getenv("SP_CLIENT_ID")
SP_CERT_THUMB  = os.getenv("SP_CERT_THUMBPRINT")
SP_CERT_PEM    = os.getenv("SP_CERT_PEM_PATH")

# ========== OpenAI ==========
from openai import OpenAI
oai = OpenAI(api_key=OPENAI_API_KEY)
EMBED_MODEL = "text-embedding-3-large"
GEN_MODEL = "gpt-4o-mini"

def embed_text(text: str) -> List[float]:
    text = text or ""
    resp = oai.embeddings.create(model=EMBED_MODEL, input=[text])
    return resp.data[0].embedding

def cosine_sim(a: List[float] | None, b: List[float] | None) -> float:
    if not a or not b:
        return 0.0
    va = np.array(a, dtype=np.float32)
    vb = np.array(b, dtype=np.float32)
    na = np.linalg.norm(va); nb = np.linalg.norm(vb)
    if na == 0 or nb == 0:
        return 0.0
    return float(np.dot(va, vb) / (na * nb))

def _coerce_json(s: str) -> dict:
    """Best-effort to parse JSON returned by an LLM."""
    if not isinstance(s, str):
        raise ValueError("Expected string")
    s = s.strip()
    if s.startswith("```"):
        first_nl = s.find("\n")
        if first_nl != -1:
            s = s[first_nl+1:]
        if s.endswith("```"):
            s = s[:-3].strip()
    start = s.find("{"); end = s.rfind("}")
    if start != -1 and end != -1 and end > start:
        s = s[start:end+1]
    s = (s.replace("\u201c", '"').replace("\u201d", '"')
           .replace("\u2018", "'").replace("\u2019", "'"))
    import re as _re
    s = _re.sub(r",\s*(\}|\])", r"\1", s)
    return json.loads(s)

def llm_structured_extract(corpus: str, content_owner: str, function: str, site: str) -> Dict[str, str]:
    """
    Extract Title/ContentSummary/Benefits from corpus via LLM,
    and attach parsed metadata (ContentOwner/Function/Site).
    """
    sys = (
        "You are an expert knowledge engineer. From the provided corpus, extract:\n"
        "1) Title: a concise representative title\n"
        "2) Content Summary: 1–2 cohesive paragraphs\n"
        "3) Benefits: 1–2 cohesive paragraphs highlighting business benefits\n"
        "Return a strict JSON object with keys: Title, ContentSummary, Benefits."
    )
    user = f"CORPUS:\n{corpus[:12000]}"
    r = oai.chat.completions.create(
        model=GEN_MODEL,
        messages=[{"role":"system","content":sys},{"role":"user","content":user}],
        temperature=0.2
    )
    text = r.choices[0].message.content
    try:
        j = _coerce_json(text)
        return {
            "Title": j.get("Title","").strip(),
            "ContentSummary": j.get("ContentSummary","").strip(),
            "Benefits": j.get("Benefits","").strip(),
            "ContentOwner": content_owner,
            "Function": function,
            "Site": site,
        }
    except Exception:
        return {
            "Title": "Untitled Knowledge Asset",
            "ContentSummary": text.strip()[:2000],
            "Benefits": "",
            "ContentOwner": content_owner,
            "Function": function,
            "Site": site,
        }

# ========== Chroma Vector DB ==========
import chromadb
try:
    client = chromadb.PersistentClient(path=VECTOR_DB_PATH)
except TypeError:
    from chromadb.config import Settings
    client = chromadb.Client(Settings(persist_directory=VECTOR_DB_PATH))

collection = client.get_or_create_collection(
    name="ai4km_assets",
    metadata={"hnsw:space":"cosine"}
)

def upsert_asset_vector(record_id: str, text: str, metadata: Dict[str, Any], vector: List[float]):
    try:
        collection.delete(ids=[record_id])
    except Exception:
        pass
    collection.add(
        ids=[record_id],
        documents=[text],
        metadatas=[metadata],
        embeddings=[vector]
    )

def search_topk(query: str, k: int = 3):
    qvec = embed_text(query)
    res = collection.query(
        query_embeddings=[qvec],
        n_results=k,
        include=["documents", "metadatas", "distances"]
    )
    ids       = (res.get("ids") or [[]])[0]
    documents = (res.get("documents") or [[]])[0]
    metadatas = (res.get("metadatas") or [[]])[0]
    distances = (res.get("distances") or [[]])[0]
    out = []
    for i in range(len(ids)):
        dist = float(distances[i]) if i < len(distances) and distances[i] is not None else None
        sim  = (1.0 - dist) if dist is not None else None
        out.append({
            "RecordId": ids[i],
            "Document": documents[i] if i < len(documents) else None,
            "Metadata": metadatas[i] if i < len(metadatas) else {},
            "Score": sim
        })
    return out

# ========== Local metadata store (CSV fallback) ==========
LOCAL_META_CSV = "vectorstore/metadata.csv"
os.makedirs("vectorstore", exist_ok=True)
if not os.path.exists(LOCAL_META_CSV):
    pd.DataFrame(columns=["RecordId","Title","ContentSummary","Benefits","ContentOwner","Function","Site"]).to_csv(LOCAL_META_CSV, index=False)

def local_meta_upsert(row: Dict[str,str]):
    df = pd.read_csv(LOCAL_META_CSV)
    df = df[df["RecordId"] != row["RecordId"]]
    df = pd.concat([df, pd.DataFrame([row])], ignore_index=True)
    df.to_csv(LOCAL_META_CSV, index=False)

def local_meta_lookup(record_ids: List[str]) -> pd.DataFrame:
    df = pd.read_csv(LOCAL_META_CSV)
    return df[df["RecordId"].isin(record_ids)].copy()

# ========== SharePoint Client (certificate app-only) ==========
def sharepoint_available() -> bool:
    needed = [SP_SITE_URL, SP_LIST_NAME, SP_CLIENT_ID, SP_TENANT, SP_CERT_THUMB, SP_CERT_PEM]
    return SHAREPOINT_ENABLED and all(needed)

def _get_sp_ctx():
    from office365.sharepoint.client_context import ClientContext
    return ClientContext(SP_SITE_URL).with_client_certificate(
        tenant=SP_TENANT,
        client_id=SP_CLIENT_ID,
        thumbprint=SP_CERT_THUMB,
        cert_path=SP_CERT_PEM
    )

def sharepoint_add_item(row: Dict[str, str]) -> bool:
    try:
        ctx = _get_sp_ctx()
        sp_list = ctx.web.lists.get_by_title(SP_LIST_NAME)
        sp_list.add_item({
            "Title":          row["Title"],
            "ContentSummary": row["ContentSummary"],
            "Benefits":       row["Benefits"],
            "ContentOwner":   row["ContentOwner"],
            "Function":       row["Function"],
            "Site":           row["Site"],
            "RecordId":       row["RecordId"],
        })
        ctx.execute_query()
        return True
    except Exception as e:
        st.warning(f"SharePoint write failed, storing locally instead. Details: {e}")
        return False

def sharepoint_fetch_by_record_id(record_id: str) -> Dict[str, Any] | None:
    try:
        ctx = _get_sp_ctx()
        sp_list = ctx.web.lists.get_by_title(SP_LIST_NAME)
        items = (
            sp_list.items
            .select(["Id","Title","RecordId","ContentSummary","Benefits","ContentOwner","Function","Site"])
            .filter(f"RecordId eq '{record_id}'")
            .top(1)
        )
        ctx.load(items); ctx.execute_query()
        if len(items) == 0:
            return None
        p = items[0].properties
        return {
            "RecordId":       p.get("RecordId"),
            "Title":          p.get("Title"),
            "ContentSummary": p.get("ContentSummary"),
            "Benefits":       p.get("Benefits"),
            "ContentOwner":   p.get("ContentOwner"),
            "Function":       p.get("Function"),
            "Site":           p.get("Site"),
            "ItemId":         p.get("Id") or p.get("ID"),
        }
    except Exception as e:
        st.warning(f"SharePoint read failed; falling back to local metadata where possible. Details: {e}")
        return None

def sharepoint_lookup_record_ids(record_ids: List[str]) -> pd.DataFrame:
    out = []
    for rid in record_ids:
        row = sharepoint_fetch_by_record_id(rid)
        if row:
            out.append(row)
    return pd.DataFrame(out) if out else pd.DataFrame(
        columns=["RecordId","Title","ContentSummary","Benefits","ContentOwner","Function","Site","ItemId"]
    )

# ========== Generation helpers ==========
def _normalize_doc_rows(df: pd.DataFrame) -> List[Dict[str, Any]]:
    """
    Produce a clean, consistent schema per doc so the LLM always sees
    Title/ContentSummary/Benefits/ContentOwner/Function/Site.
    Prefers *_sp (SharePoint-enriched) columns when available.
    """
    records = df.to_dict(orient="records")
    out = []
    for r in records:
        out.append({
            "RecordId":       r.get("RecordId", ""),
            "Title":          (r.get("Title") or r.get("Title_sp") or "").strip(),
            "ContentSummary": (r.get("ContentSummary") or "").strip(),
            "Benefits":       (r.get("Benefits") or "").strip(),
            "ContentOwner":   (r.get("ContentOwner") or r.get("ContentOwner_sp") or "").strip(),
            "Function":       (r.get("Function") or r.get("Function_sp") or "").strip(),
            "Site":           (r.get("Site") or r.get("Site_sp") or "").strip(),
        })
    return out

def _maybe_answer_from_metadata(user_query: str, docs: List[Dict[str, Any]]) -> str | None:
    q = user_query.lower()

    def unique_nonempty(field: str) -> list[str]:
        vals = [(d.get(field) or "").strip() for d in docs]
        vals = [v for v in vals if v]
        return sorted(set(vals))

    # Site questions
    if "which site" in q or "what site" in q or "site implemented" in q or "site" in q:
        sites = unique_nonempty("Site")
        if sites:
            if len(sites) == 1:
                return f"The improvement was implemented at **{sites[0]}**."
            else:
                return "The improvement appears in these sites: " + ", ".join(f"**{s}**" for s in sites) + "."

    # Owner questions
    if "owner" in q or "content owner" in q or "who owns" in q:
        owners = unique_nonempty("ContentOwner")
        if owners:
            if len(owners) == 1:
                return f"The content owner is **{owners[0]}**."
            else:
                return "Relevant content owners: " + ", ".join(f"**{o}**" for o in owners) + "."

    # Function/department questions
    if "function" in q or "department" in q:
        funcs = unique_nonempty("Function")
        if funcs:
            if len(funcs) == 1:
                return f"The function/department is **{funcs[0]}**."
            else:
                return "Relevant functions: " + ", ".join(f"**{f}**" for f in funcs) + "."

    return None

def generate_answer_with_history(query: str,
                                 docs: List[Dict[str, Any]],
                                 history: List[Dict[str, str]] | None = None,
                                 max_history_turns: int = 6) -> str:
    """
    Answer with LLM using retrieved docs + recent chat history.
    Includes metadata fields (Site, ContentOwner, Function) in context.
    """
    meta_answer = _maybe_answer_from_metadata(query, docs)
    if meta_answer:
        return meta_answer

    if not docs:
        base_context = "No relevant documents found to answer your question."
    else:
        parts = []
        for i, d in enumerate(docs, start=1):
            parts.append(
                f"[Doc {i}] RecordId: {d.get('RecordId','')}\n"
                f"Title: {d.get('Title','')}\n"
                f"Site: {d.get('Site','')}\n"
                f"Function: {d.get('Function','')}\n"
                f"ContentOwner: {d.get('ContentOwner','')}\n"
                f"Summary: {d.get('ContentSummary','')}\n"
                f"Benefits: {d.get('Benefits','')}\n"
            )
        base_context = "\n\n".join(parts)

    history = history or []
    clipped = history[-max_history_turns:]

    system_prompt = (
        "You are an expert enterprise assistant.\n"
        "Use ONLY the provided context from knowledge assets. If the context is insufficient, say so.\n"
        "When the user asks about fields like Site, Function, or ContentOwner, extract them directly from the context.\n"
        "Be concise and factual (1–2 short paragraphs). Do not invent details."
    )

    msgs = [{"role": "system", "content": system_prompt}]
    msgs.append({"role": "assistant", "content": f"Context for answering:\n{base_context}"})
    for m in clipped:
        msgs.append({"role": m["role"], "content": m["content"]})
    msgs.append({"role": "user", "content": query})

    r = oai.chat.completions.create(model=GEN_MODEL, messages=msgs, temperature=0.2)
    return r.choices[0].message.content.strip()

def suggest_followups(user_query: str, docs: List[Dict[str, Any]], answer: str, max_suggestions: int = 3) -> list[str]:
    """
    Proposes grounded, helpful next questions based on the current query, the retrieved docs' metadata,
    and the assistant's answer. Returns <= max_suggestions items.
    """
    if docs:
        meta_lines = []
        for i, d in enumerate(docs[:5], start=1):
            meta_lines.append(
                f"[Doc {i}] Title={d.get('Title','')[:140]} | Site={d.get('Site','')} | "
                f"Function={d.get('Function','')} | Owner={d.get('ContentOwner','')}"
            )
        doc_meta = "\n".join(meta_lines)
    else:
        doc_meta = "No documents."

    sys = (
        "You are a helpful assistant that proposes follow-up questions.\n"
        "Rules:\n"
        "- Only suggest questions that can likely be answered from the provided documents' metadata or summaries.\n"
        "- Keep each question short, specific, and non-redundant.\n"
        "- Avoid generic tutorials or policy questions unless they clearly relate to the docs.\n"
        "- Do NOT answer; just return a JSON object: {\"suggestions\": [\"...\", \"...\", \"...\"]}."
    )
    user = (
        f"UserQuery:\n{user_query}\n\n"
        f"Docs:\n{doc_meta}\n\n"
        f"AssistantAnswer:\n{answer[:1200]}"
    )

    try:
        r = oai.chat.completions.create(
            model=GEN_MODEL,
            messages=[{"role":"system","content":sys}, {"role":"user","content":user}],
            temperature=0.2
        )
        j = _coerce_json(r.choices[0].message.content)
        out = [s.strip() for s in (j.get("suggestions") or []) if isinstance(s, str) and s.strip()]
        uniq = []
        for s in out:
            if s not in uniq:
                uniq.append(s)
        return uniq[:max_suggestions] if uniq else []
    except Exception:
        return []

def run_rag_turn(user_query: str,
                 prev_query_vec: List[float] | None,
                 context_buffer: List[Dict[str, Any]],
                 topic_shift_threshold: float = 0.60
                 ) -> tuple[str, pd.DataFrame, bytes, List[float], List[Dict[str, Any]], list[str]]:
    """
    One conversational RAG turn with topic-shift detection and context aggregation.
    Returns: (answer, out_df, excel_bytes, curr_query_vec, updated_context_buffer, suggestions)
    """
    curr_vec = embed_text(user_query)
    if prev_query_vec is None or cosine_sim(prev_query_vec, curr_vec) < topic_shift_threshold:
        context_buffer = []

    hits = search_topk(user_query, k=3)
    if not hits:
        empty_df = pd.DataFrame(columns=["RecordId","Similarity","Title","ContentOwner","Function","Site"])
        return ("I couldn't find relevant knowledge assets for that query.", empty_df, b"", curr_vec, context_buffer, [])

    rows = []
    for h in hits:
        m = h["Metadata"]
        rows.append({
            "RecordId": h["RecordId"],
            "Similarity": round(h["Score"], 4) if h["Score"] is not None else None,
            "Title": m.get("Title",""),
            "ContentOwner": m.get("ContentOwner",""),
            "Function": m.get("Function",""),
            "Site": m.get("Site","")
        })
    df = pd.DataFrame(rows)

    rec_ids = [r["RecordId"] for r in rows]
    if sharepoint_available():
        sp_extra = sharepoint_lookup_record_ids(rec_ids)
        sp_extra = sp_extra.rename(columns={
            "Title": "Title_sp",
            "ContentOwner": "ContentOwner_sp",
            "Function": "Function_sp",
            "Site": "Site_sp"
        })
        sp_cols = ["RecordId", "ContentSummary", "Benefits", "Title_sp", "ContentOwner_sp", "Function_sp", "Site_sp"]
        sp_extra = sp_extra[sp_cols] if not sp_extra.empty else sp_extra
        out = df.merge(sp_extra, on="RecordId", how="left")
    else:
        local_extra = local_meta_lookup(rec_ids)
        out = df.merge(local_extra, on="RecordId", how="left", suffixes=("",""))

    norm_docs_new = _normalize_doc_rows(out)

    combined = norm_docs_new + (context_buffer or [])
    dedup = {}
    for d in combined:
        dedup[d.get("RecordId")] = d
    updated_buffer = list(dedup.values())[:6]

    answer = generate_answer_with_history(user_query, updated_buffer, st.session_state.chat_msgs)
    suggestions = suggest_followups(user_query, updated_buffer, answer, max_suggestions=3)

    buffer = io.BytesIO()
    with pd.ExcelWriter(buffer, engine="openpyxl") as writer:
        out.to_excel(writer, index=False, sheet_name="Matches")
    excel_bytes = buffer.getvalue()

    return answer, out, excel_bytes, curr_vec, updated_buffer, suggestions

def _append_assistant_turn(answer: str, out_df: pd.DataFrame, excel_bytes: bytes, suggestions: list[str]):
    """Append a single assistant message to chat state (no UI rendering here)."""
    st.session_state.chat_msgs.append({
        "role": "assistant",
        "content": answer,
        "matches_df": out_df,
        "excel_bytes": excel_bytes,
        "suggestions": suggestions,
    })

def _do_query_and_append(query: str):
    """Run one RAG turn for `query`, update state, do not render. Rendering happens in the history loop."""
    st.session_state.chat_msgs.append({"role": "user", "content": query})
    answer, out_df, excel_bytes, curr_vec, updated_buffer, suggestions = run_rag_turn(
        query,
        prev_query_vec=st.session_state.prev_query_vec,
        context_buffer=st.session_state.context_buffer
    )
    st.session_state.prev_query_vec = curr_vec
    st.session_state.context_buffer = updated_buffer
    st.session_state.last_excel_bytes = excel_bytes
    _append_assistant_turn(answer, out_df, excel_bytes, suggestions)

# ========== File parsers ==========
from PyPDF2 import PdfReader
from docx import Document
import tempfile
from pptx import Presentation

def read_pdf(file) -> str:
    with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tmp:
        file.seek(0); tmp.write(file.read()); tmp.flush()
        reader = PdfReader(tmp.name)
        texts = []
        for page in reader.pages:
            try:
                texts.append(page.extract_text() or "")
            except Exception:
                pass
        return "\n".join(texts)

def read_docx(file) -> str:
    with tempfile.NamedTemporaryFile(delete=False, suffix=".docx") as tmp:
        file.seek(0); tmp.write(file.read()); tmp.flush()
        doc = Document(tmp.name)
        return "\n".join([p.text for p in doc.paragraphs])

def read_xlsx(file) -> str:
    with tempfile.NamedTemporaryFile(delete=False, suffix=".xlsx") as tmp:
        file.seek(0); tmp.write(file.read()); tmp.flush()
        df = pd.read_excel(tmp.name, sheet_name=None)
        parts = []
        for name, sdf in df.items():
            parts.append(f"## Sheet: {name}")
            parts.append(sdf.to_string(index=False))
        return "\n".join(parts)

def read_txt(file) -> str:
    file.seek(0)
    return file.read().decode("utf-8", errors="ignore")

def read_pptx(file) -> str:
    with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp:
        file.seek(0); tmp.write(file.read()); tmp.flush()
        path = tmp.name
    prs = Presentation(path)
    parts = []
    for idx, slide in enumerate(prs.slides, start=1):
        parts.append(f"## Slide {idx}")
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = "".join(run.text for run in para.runs).strip()
                    if t:
                        parts.append(t)
            if getattr(shape, "has_table", False) and shape.has_table:
                tbl = shape.table
                for row in tbl.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    row_text = " | ".join([c for c in cells if c])
                    if row_text:
                        parts.append(row_text)
        try:
            if getattr(slide, "has_notes_slide", False) and slide.has_notes_slide:
                ns = slide.notes_slide
                note_chunks = []
                try:
                    if ns and ns.notes_text_frame and ns.notes_text_frame.text:
                        note_chunks.append(ns.notes_text_frame.text.strip())
                except Exception:
                    pass
                try:
                    for shp in getattr(ns, "shapes", []):
                        if getattr(shp, "has_text_frame", False) and shp.has_text_frame:
                            for para in shp.text_frame.paragraphs:
                                t = "".join(run.text for run in para.runs).strip()
                                if t:
                                    note_chunks.append(t)
                except Exception:
                    pass
                if note_chunks:
                    parts.append("### Notes")
                    parts.extend(note_chunks)
        except Exception:
            pass
    return "\n".join([p for p in parts if p.strip()])

def consolidate_files(files) -> str:
    corpus_parts = []
    for f in files:
        name = f.name.lower()
        if name.endswith(".pdf"):
            corpus_parts.append(read_pdf(f))
        elif name.endswith(".docx"):
            corpus_parts.append(read_docx(f))
        elif name.endswith(".xlsx") or name.endswith(".xls"):
            corpus_parts.append(read_xlsx(f))
        elif name.endswith(".txt"):
            corpus_parts.append(read_txt(f))
        elif name.endswith(".pptx"):
            corpus_parts.append(read_pptx(f))
        else:
            corpus_parts.append(read_txt(f))  # naive fallback
    return "\n\n".join([p for p in corpus_parts if p.strip()])

# ========== Validators & Metadata Extraction ==========
AMGEN_EMAIL_RE = re.compile(r"^[A-Za-z0-9._%+-]+@amgen\.com$")
SITE_RE = re.compile(r"^[A-Za-z]{3}$")

def to_title_case_words(s: str) -> str:
    return " ".join([w.capitalize() for w in re.split(r"\s+", (s or "").strip()) if w])

META_PATTERNS = {
    "email": re.compile(r"(?im)\bemail\s*[:\-–]\s*([A-Za-z0-9._%+\-]+@amgen\.com)\b"),
    "site": re.compile(r"(?im)\bsite\s*[:\-–]\s*([A-Za-z]{2,10})\b"),
    "function": re.compile(r"(?im)\bfunction\s*[:\-–]\s*(.+)$"),
}

def extract_metadata_from_text(corpus: str) -> Dict[str, str]:
    email = site = function = ""
    m = META_PATTERNS["email"].search(corpus)
    if m:
        email = m.group(1).strip().lower()
    m = META_PATTERNS["site"].search(corpus)
    if m:
        site_candidate = m.group(1).strip().upper()
        site_candidate = re.sub(r"[^A-Z]", "", site_candidate)[:3]
        site = site_candidate
    m = META_PATTERNS["function"].search(corpus)
    if m:
        function_candidate = m.group(1).strip().rstrip(".,;: ")
        function = to_title_case_words(function_candidate)
    return {"email": email, "site": site, "function": function}

# ========== UI ==========
st.set_page_config(page_title="AI4KM RAG Pipeline", page_icon="💡", layout="wide")

# ---------- Session keys ----------
def _reset_ingest_state():
    st.session_state.pop("ing_meta", None)
    st.session_state.pop("ing_corpus", None)
    st.session_state.pop("ing_ready", None)
    st.session_state.pop("ing_last_record_id", None)
    st.session_state.pop("ing_parsed_meta", None)

# init keys once (ingestion)
for k in ["ing_meta", "ing_corpus", "ing_ready", "ing_last_record_id", "ing_parsed_meta", "ing_flow_done"]:
    if k not in st.session_state:
        st.session_state[k] = None

# chat state (search)
for k in ["chat_msgs", "last_excel_bytes", "prev_query_vec", "context_buffer", "pending_auto_query"]:
    if k not in st.session_state:
        st.session_state[k] = None
if st.session_state.pending_auto_query is None:
    st.session_state.pending_auto_query = ""
if st.session_state.chat_msgs is None:
    st.session_state.chat_msgs = []
if st.session_state.last_excel_bytes is None:
    st.session_state.last_excel_bytes = b""
if st.session_state.prev_query_vec is None:
    st.session_state.prev_query_vec = None
if st.session_state.context_buffer is None:
    st.session_state.context_buffer = []

# ============= Global Styles and Animated Background =============
st.markdown("""
<style>
/* ----- App background ----- */
html, body, [data-testid="stAppViewContainer"] { height: 100%; margin: 0; }
[data-testid="stAppViewContainer"]{
  position: relative;
  background: radial-gradient(circle at center,
    #e6f2ff 0%, #cfe7ff 26%, #84bdfd 48%, #3c87d3 70%, #183160 90%, #09152d 100%);
}
header[data-testid="stHeader"], .block-container{ position:relative; z-index:1; }

/* ----- Clouds (base layer) ----- */
.km-sky{
  position: fixed; inset: 0; overflow: hidden;
  pointer-events: none; z-index:0;
}
.km-cloud{
  position: absolute; left: -40vw;
  width: 36vw; max-width: 520px; aspect-ratio: 16/9;
  opacity:.93; filter: drop-shadow(0 12px 16px rgba(0,0,0,.08));
  animation: cloudDrift var(--dur,22s) linear infinite;
  animation-delay: var(--delay,0s);
  background:
    radial-gradient(35% 45% at 20% 65%, rgba(255,255,255,.98) 0 60%, transparent 61%),
    radial-gradient(40% 50% at 40% 55%, rgba(255,255,255,.98) 0 60%, transparent 61%),
    radial-gradient(45% 52% at 60% 58%, rgba(255,255,255,.98) 0 60%, transparent 61%),
    radial-gradient(34% 44% at 75% 62%, rgba(255,255,255,.98) 0 60%, transparent 61%),
    radial-gradient(30% 40% at 50% 40%, rgba(255,255,255,.98) 0 60%, transparent 61%),
    radial-gradient(60% 55% at 50% 70%, rgba(255,255,255,.96) 0 62%, transparent 63%);
}
.km-cloud.c1{top:18%; --scale:.95; --dur:18s; --delay:-8s;}
.km-cloud.c2{top:32%; --scale:.85; --dur:20s; --delay:-14s;}
.km-cloud.c3{top:46%; --scale:1.05; --dur:19s; --delay:-4s;}
.km-cloud.c4{top:60%; --scale:.92; --dur:21s; --delay:-20s;}
.km-cloud.c5{top:72%; --scale:1.15; --dur:22s; --delay:-10s;}
@keyframes cloudDrift{
  from{transform:translateX(0) translateY(-50%) scale(var(--scale,1));}
  to  {transform:translateX(200vw) translateY(-50%) scale(var(--scale,1));}
}

/* ===== Card wrappers (reliable, visible colors) ===== */
.km-card{
  border-radius: 22px;
  padding: 18px 16px 22px;
  margin: 6px 8px 10px;
  box-shadow: 0 14px 30px rgba(0,0,0,.18);
  position: relative;
  z-index: 2; /* above clouds */
}

/* Bright SOLID backgrounds (no gradients to avoid blending with sky) */
.km-ingest { background: #ff7a00; }   /* bright orange */
.km-search { background: #16c172; }   /* bright green  */

/* Inner white panels so inputs are readable; semi-opaque so color frame shows */
.km-pane{
  background: rgba(255,255,255,.96);
  border-radius:16px;
  padding:12px 14px;
  box-shadow: inset 0 1px 0 rgba(255,255,255,.6);
  margin:12px 0;
}

/* Chat bubbles */
.km-chat{
  /* Make the chat pane tall and scrollable so the page itself doesn't grow */
  height: calc(100vh - 360px);   /* adjust if your header/subheaders change */
  min-height: 300px;
  max-height: 40vh;              /* safety cap on very tall screens */
  overflow-y: auto;
  padding-right: 6px;
  scroll-behavior: smooth;
  overscroll-behavior: contain;  /* avoid "rubber band" scroll affecting page */
}
.km-msg{ margin:8px 0; padding:10px 12px; border-radius:12px; background:rgba(255,255,255,.92); }
.km-msg.user{ border-left:4px solid #0D6EFD; }
.km-msg.assistant{ border-left:4px solid #6633ff; }

/* Keep widgets from stretching edge-to-edge */
[data-testid="stTextInputRoot"],
[data-testid="stFileUploaderDropzone"],
[data-testid="stDataFrameContainer"],
[data-testid="stDownloadButton"]{ max-width:880px; }
.stButton button, [data-testid="stFormSubmitButton"] button{ width:auto; padding:.55rem 1rem; }

/* ---- Bold, bigger text helpers ---- */
.km-lead{
  font-weight: 700;
  font-size: 1.05rem;   /* slightly larger than normal */
  color: #071a33;
  margin: 4px 0 12px;
}
.km-label{
  font-weight: 700;
  font-size: 1rem;
  color: #071a33;
  margin: 6px 0 6px;
  display: block;
}
.km-note{
  font-weight: 600;
  font-size: .95rem;
  color: #071a33;
}
</style>

<!-- clouds -->
<div class="km-sky">
  <div class="km-cloud c1"></div>
  <div class="km-cloud c2"></div>
  <div class="km-cloud c3"></div>
  <div class="km-cloud c4"></div>
  <div class="km-cloud c5"></div>
</div>
""", unsafe_allow_html=True)

st.title("AI4KM (AI for Knowledge Marketplace)")
left, right = st.columns([1, 1], gap="large")

# ---------------- LEFT: DATA INGESTION ----------------
with left:
    st.markdown('<div class="km-card km-ingest">', unsafe_allow_html=True)

    st.subheader("Data Ingestion")
    #st.caption("Submit new knowledge assets. I’ll extract Title, Summary, Benefits and index them.")
    st.markdown('<div class="km-lead">Submit new knowledge assets. I’ll extract <b>Title</b>, <b>Summary</b>, <b>Benefits</b> and index them.</div>', unsafe_allow_html=True)

    # Success panel (after SharePoint ok)
    show_success_panel = (st.session_state.get("ing_flow_done") == "sp_ok")
    if show_success_panel:
        rid = st.session_state.get("ing_last_record_id", "")
        st.success(f"✅ Submitted to SharePoint successfully. RecordId: {rid}")
        if st.button("➕ Ingest another?"):
            st.session_state.ing_flow_done = None
            st.session_state.ing_last_record_id = None
            st.rerun()

    if not show_success_panel:
        if not st.session_state.ing_ready:
            with st.form("ingestion_form_card", clear_on_submit=False):
                st.markdown('<div class="km-pane">', unsafe_allow_html=True)
                st.markdown('<label class="km-label">Attach reference file(s) (drag & drop)</label>', unsafe_allow_html=True)                
                files = st.file_uploader(
                    "",
                    accept_multiple_files=True,
                    type=["pdf","docx","xlsx","xls","txt","pptx"],
                    label_visibility="collapsed"
                )
                st.markdown('</div>', unsafe_allow_html=True)
                submit_ing = st.form_submit_button("Process")

            if submit_ing:
                if not files:
                    st.error("Please attach at least one file.")
                else:
                    with st.spinner("Reading and consolidating files…"):
                        corpus = consolidate_files(files)
                        if not corpus.strip():
                            st.error("No readable text found in the uploaded files.")
                        else:
                            with st.expander("🪲 Debug: Show consolidated corpus (first 10,000 chars)"):
                                st.caption(f"Corpus length: {len(corpus):,} characters")
                                st.text_area("corpus", value=corpus[:10000], height=280)
                            st.download_button(
                                "⬇️ Download full corpus (debug)",
                                data=corpus.encode("utf-8", errors="ignore"),
                                file_name="debug_corpus.txt",
                                mime="text/plain",
                            )
                            parsed = extract_metadata_from_text(corpus)
                            errors = []
                            if not parsed["email"] or not AMGEN_EMAIL_RE.match(parsed["email"]):
                                errors.append("Email (must be an @amgen.com address)")
                            if not parsed["site"] or not SITE_RE.match(parsed["site"]):
                                errors.append("Site (3-letter code, e.g., ASM)")
                            if not parsed["function"]:
                                errors.append("Function/Department")
                            if errors:
                                st.error(
                                    "I couldn't find valid values for:\n\n- "
                                    + "\n- ".join(errors)
                                    + "\n\nPlease make sure your file contains lines like:\n"
                                    "`Email: jdoe@amgen.com`, `Site: ASM`, `Function: Digital Technology & Innovation`"
                                )
                            else:
                                with st.spinner("Generating Title, Content Summary and Benefits via LLM…"):
                                    meta = llm_structured_extract(
                                        corpus,
                                        content_owner=parsed["email"],
                                        function=parsed["function"],
                                        site=parsed["site"],
                                    )
                                st.session_state.ing_corpus = corpus
                                st.session_state.ing_meta = meta
                                st.session_state.ing_parsed_meta = parsed
                                st.session_state.ing_ready = True
                                st.rerun()
        else:
            meta = st.session_state.ing_meta or {}
            parsed = st.session_state.ing_parsed_meta or {}

            st.success("Draft extracted. Please review:")
            st.markdown('<div class="km-pane">', unsafe_allow_html=True)
            st.write(f"**Title**: {meta.get('Title','')}")
            st.write("**Content Summary**"); st.write(meta.get("ContentSummary",""))
            st.write("**Benefits**"); st.write(meta.get("Benefits",""))
            st.markdown('</div>', unsafe_allow_html=True)

            with st.expander("Parsed Metadata (from your files)"):
                st.markdown(
                    f"- **Email**: `{parsed.get('email','')}`\n"
                    f"- **Site**: `{parsed.get('site','')}`\n"
                    f"- **Function**: `{parsed.get('function','')}`"
                )

            st.info("Do you want to submit a new knowledge asset?")
            coly, coln = st.columns(2)
            yes = coly.button("Yes — Submit", key="yes_submit_card")
            no  = coln.button("No — Cancel", key="no_cancel_card")

            if yes:
                record_id = str(uuid.uuid4())
                corpus = st.session_state.ing_corpus or ""
                with st.spinner("Creating embedding and saving to vector DB…"):
                    doc_for_vector = corpus[:200000] if len(corpus) > 200000 else corpus
                    vec = embed_text(doc_for_vector)
                    upsert_asset_vector(
                        record_id=record_id,
                        text=doc_for_vector,
                        metadata={
                            "Title": meta.get("Title",""),
                            "ContentOwner": meta.get("ContentOwner",""),
                            "Function": meta.get("Function",""),
                            "Site": meta.get("Site","")
                        },
                        vector=vec
                    )
                row = {"RecordId": record_id, **meta}
                sp_ok = False
                if sharepoint_available():
                    st.write("Attempting SharePoint write…")
                    if sharepoint_add_item(row):
                        sp_ok = True
                if not sp_ok:
                    local_meta_upsert(row)
                    st.info("Saved locally (CSV). You can switch to SharePoint later by setting SHAREPOINT_ENABLED=true.")

                # clear draft state after submission attempt
                st.session_state.ing_meta = None
                st.session_state.ing_corpus = None
                st.session_state.ing_ready = None
                st.session_state.ing_parsed_meta = None

                # post-submit outcome
                st.session_state.ing_last_record_id = record_id
                st.session_state.ing_flow_done = "sp_ok" if sp_ok else None

                st.rerun()

            if no:
                st.session_state.ing_meta = None
                st.session_state.ing_corpus = None
                st.session_state.ing_ready = None
                st.session_state.ing_parsed_meta = None
                st.rerun()

    st.markdown('</div>', unsafe_allow_html=True)  # close orange card

# ---------------- RIGHT: INTELLIGENT SEARCH ----------------
with right:
    st.markdown('<div class="km-card km-search">', unsafe_allow_html=True)

    st.subheader("Intelligent Search")
    st.markdown(
        '<div class="km-lead">Ask about knowledge assets. I’ll <b>retrieve</b>, '
        '<b>ground answers</b>, and <b>suggest follow-ups</b>.</div>',
        unsafe_allow_html=True
    )

    # --- Role alignment overrides (User -> right, Assistant -> left) ---
    st.markdown("""
    <style>
      .km-msg{ max-width:92%; }
      .km-msg.user{
        margin-left:auto; text-align:right;
        border-left:none; border-right:4px solid #0D6EFD;
      }
      .km-msg.assistant{
        margin-right:auto; text-align:left;
        border-right:none; border-left:4px solid #6633ff;
      }
    </style>
    """, unsafe_allow_html=True)

    # 1) Handle any auto-question first (from suggestion chips)
    if st.session_state.pending_auto_query:
        auto_q = st.session_state.pending_auto_query
        st.session_state.pending_auto_query = ""
        with st.spinner("Searching and generating…"):
            _do_query_and_append(auto_q)
            st.rerun()  # ensure render in the same run

    # 2) Bottom-pinned composer FIRST (so new turns are processed before history renders)
    chat_prompt = st.chat_input("Type your query…")
    if chat_prompt:
        with st.spinner("Searching and generating…"):
            _do_query_and_append(chat_prompt)
            st.rerun()  # render freshly-added turn immediately

    # 3) Scrollable chat box: render ONLY message bubbles inside it
    bubble_html_parts = []
    msgs = st.session_state.chat_msgs or []
    if not msgs:
        bubble_html_parts.append('<div class="km-note">Your conversation will appear here.</div>')
    else:
        for m in msgs:
            role = m.get("role", "assistant")
            css_role = "user" if role == "user" else "assistant"
            txt = m.get("content", "")
            bubble_html_parts.append(
                f'<div class="km-msg {css_role}"><strong>{role.title()}:</strong><br>{txt}</div>'
            )

    st.markdown(
        f'<div class="km-pane km-chat">{"".join(bubble_html_parts)}</div>',
        unsafe_allow_html=True
    )

    # 4) Render ONLY the latest assistant turn's attachments (DF + follow-ups)
    last_asst_idx = None
    for i in range(len(msgs) - 1, -1, -1):
        if msgs[i].get("role") == "assistant":
            last_asst_idx = i
            break

    if last_asst_idx is not None:
        last_asst = msgs[last_asst_idx]

        # Latest results table (Update 1)
        if isinstance(last_asst.get("matches_df"), pd.DataFrame) and not last_asst["matches_df"].empty:
            st.dataframe(last_asst["matches_df"], use_container_width=True)

        # Latest follow-up suggestions (Update 2)
        if last_asst.get("suggestions"):
            st.write("**You could also ask:**")
            sug_cols = st.columns(len(last_asst["suggestions"]))
            for i, s in enumerate(last_asst["suggestions"]):
                if sug_cols[i].button(s, key=f"sugg_turn_latest_{i}"):
                    st.session_state.pending_auto_query = s
                    st.rerun()

    # 5) Controls row (Update 3 already applied by removing per-turn downloads)
    c1, c2 = st.columns(2)
    with c1:
        if st.button("🧹 Clear conversation", use_container_width=True):
            st.session_state.chat_msgs = []
            st.session_state.last_excel_bytes = b""
            st.session_state.prev_query_vec = None
            st.session_state.context_buffer = []
            st.session_state.pending_auto_query = ""
            st.rerun()
    with c2:
        if st.session_state.last_excel_bytes:
            st.download_button(
                "⬇️ Download most recent Top-3 (Excel)",
                data=st.session_state.last_excel_bytes,
                file_name="ai4km_search_results.xlsx",
                mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                use_container_width=True,
                key="dl_latest_card"
            )

    st.markdown('</div>', unsafe_allow_html=True)  # close green card
